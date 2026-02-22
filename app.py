"""Simple assignment grader using Azure OpenAI."""

import os
import json
import re
import tempfile
from flask import Flask, render_template, request, jsonify
from werkzeug.utils import secure_filename
import nbformat
from docx import Document
from openai import AzureOpenAI
from dotenv import load_dotenv

try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover
    PdfReader = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None

try:
    from openpyxl import load_workbook
except ImportError:  # pragma: no cover
    load_workbook = None

load_dotenv()

# Handle deployment at /grader-gpt subdirectory
# When deployed, nginx will strip /grader-gpt prefix before forwarding
app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024
app.config['UPLOAD_FOLDER'] = tempfile.mkdtemp()

# Application root - will be /grader-gpt when deployed
APPLICATION_ROOT = os.getenv('APPLICATION_ROOT', '/')


def _read_as_plain_text(file_path):
    """Best-effort UTF text loader for generic text files."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_text(file_path):
    """Extract text from common assignment document formats."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == '.ipynb':
        nb = nbformat.read(file_path, as_version=4)
        text = []
        for cell in nb.cells:
            text.append(cell.get('source', ''))
        return '\n\n'.join(text)
    if ext == '.docx':
        doc = Document(file_path)
        text = []
        for para in doc.paragraphs:
            text.append(para.text)
        return '\n\n'.join(text)

    if ext == '.pdf':
        if not PdfReader:
            return _read_as_plain_text(file_path)
        reader = PdfReader(file_path)
        pages = [(page.extract_text() or "") for page in reader.pages]
        return "\n\n".join(pages)

    if ext == '.pptx':
        if not Presentation:
            return _read_as_plain_text(file_path)
        presentation = Presentation(file_path)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
        return "\n\n".join(text)

    if ext in {'.xlsx', '.xlsm'}:
        if not load_workbook:
            return _read_as_plain_text(file_path)
        wb = load_workbook(file_path, data_only=True)
        text = []
        for sheet in wb.worksheets:
            text.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(v) for v in row if v is not None]
                if row_values:
                    text.append(" | ".join(row_values))
        return "\n".join(text)

    text_like_extensions = {
        '.txt', '.md', '.csv', '.json', '.jsonl', '.yaml', '.yml', '.xml',
        '.html', '.htm', '.py', '.js', '.ts', '.java', '.c', '.cpp', '.h',
        '.sql', '.r', '.tex', '.sh', '.bat', '.log', '.ini', '.cfg', '.toml',
        '.rtf', '.doc'
    }

    if ext in text_like_extensions:
        return _read_as_plain_text(file_path)

    # Last-resort fallback: attempt text decode so unknown formats still get parsed if possible.
    return _read_as_plain_text(file_path)


def get_azure_client():
    """Initialize Azure OpenAI client."""
    api_key = os.getenv('AZURE_OPENAI_API_KEY')
    endpoint = os.getenv('AZURE_OPENAI_ENDPOINT')
    api_version = os.getenv('AZURE_OPENAI_API_VERSION', '2024-12-01-preview')
    
    if not api_key or not endpoint:
        raise ValueError("AZURE_OPENAI_API_KEY and AZURE_OPENAI_ENDPOINT must be set")
    
    # Ensure endpoint doesn't have trailing slash
    endpoint = endpoint.rstrip('/')
    
    print(f"Connecting to Azure OpenAI:")
    print(f"  Endpoint: {endpoint}")
    print(f"  API Version: {api_version}")
    
    return AzureOpenAI(
        api_key=api_key,
        api_version=api_version,
        azure_endpoint=endpoint
    )


@app.route('/health')
def health():
    """Health check endpoint."""
    return jsonify({'status': 'healthy'}), 200


@app.route('/')
def index():
    """Main page."""
    has_key = bool(os.getenv('AZURE_OPENAI_API_KEY'))
    return render_template('index.html', has_key=has_key)


@app.route('/grade', methods=['POST'])
def grade():
    """Grade submissions."""
    try:
        # Get files
        base_file = request.files.get('base_file')
        student_files = request.files.getlist('student_files')
        
        if not base_file or not student_files:
            return jsonify({'error': 'Missing files'}), 400
        
        # Save base file
        base_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(base_file.filename))
        base_file.save(base_path)
        base_content = extract_text(base_path)
        
        # Initialize Azure client
        client = get_azure_client()
        model = os.getenv('AZURE_OPENAI_DEPLOYMENT_NAME', 'gpt-4')
        
        print(f"Using deployment: {model}")
        
        # Grade each student file
        results = []
        for student_file in student_files:
            if not student_file.filename:
                continue
            
            student_path = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(student_file.filename))
            student_file.save(student_path)
            student_content = extract_text(student_path)
            
            # Grade with Azure OpenAI
            prompt = f"""You are grading a student's submission against a reference solution.

IMPORTANT GRADING GUIDELINES:
- Do NOT penalize minor style differences or variable naming.
- Partial credit is encouraged when the approach is reasonable.
- Focus on correctness, logic, and completeness over exact matching.

BASE SOLUTION (for reference):
------------------------------
{base_content}

STUDENT SUBMISSION:
-------------------
{student_content}

EVALUATION CRITERIA:
1. Correctness of core logic and results
2. Completeness of the solution
3. Reasonable handling of edge cases
4. Clarity of explanations or comments (if applicable)

SCORING RULES:
- Start from 100 points.
- Deduct points ONLY for clear mistakes or missing components.
- Small issues should incur small deductions (1–5 points).
- Larger conceptual errors may incur larger deductions.
- Do not invent issues if the solution is acceptable.

OUTPUT FORMAT:
Respond ONLY in valid JSON with the following structure:

{{
  "score": xx,
  "feedback": "<brief, constructive summary of strengths and weaknesses>",
  "question_feedback": [
    {{"question_number": "1", "reasoning": " "}},
    {{"question_number": "2", "reasoning": " "}}
  ],
  "deductions": [
    {{"issue": " ", "points": x, "section": " "}},
    {{"issue": " ", "points": x, "section": " "}}
  ]
}}

If there are no meaningful issues, return empty question_feedback and deductions lists."""

            # Prepare API request details for debugging
            system_message = "You are a grading assistant. Always respond with valid JSON."
            api_request = {
                "model": model,
                "messages": [
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": prompt}
                ],
                "response_format": {"type": "json_object"}
            }
            
            try:
                response = client.chat.completions.create(**api_request)
            except Exception as e:
                error_msg = str(e)
                if "DeploymentNotFound" in error_msg:
                    return jsonify({
                        'error': f'Deployment "{model}" not found. Check AZURE_OPENAI_DEPLOYMENT_NAME in .env file. Available deployments can be found in Azure Portal.'
                    }), 400
                raise
            
            raw_response = response.choices[0].message.content
            result = json.loads(raw_response)
            
            # Collect debug information
            debug_info = {
                "api_request": {
                    "model": model,
                    "endpoint": os.getenv('AZURE_OPENAI_ENDPOINT'),
                    "api_version": os.getenv('AZURE_OPENAI_API_VERSION', '2024-12-01-preview'),
                    "system_message": system_message,
                    "user_prompt": prompt,
                    "prompt_length": len(prompt),
                    "temperature": 0.3
                },
                "api_response": {
                    "model_used": response.model if hasattr(response, 'model') else model,
                    "finish_reason": response.choices[0].finish_reason,
                    "raw_response": raw_response,
                    "usage": {
                        "prompt_tokens": response.usage.prompt_tokens if hasattr(response, 'usage') else None,
                        "completion_tokens": response.usage.completion_tokens if hasattr(response, 'usage') else None,
                        "total_tokens": response.usage.total_tokens if hasattr(response, 'usage') else None
                    }
                },
                "parsed_result": result
            }
            
            # Try to extract student name/ID from content
            filename = student_file.filename
            student_name = filename.split('.')[0] if '.' in filename else filename
            student_id = None
            
            # Look for student info in content (simple extraction)
            content_lower = student_content.lower()
            if 'netid' in content_lower or 'student id' in content_lower:
                netid_match = re.search(r'(?:netid|student\s+id)\s*[:=]\s*(\S+)', student_content, re.IGNORECASE)
                if netid_match:
                    student_id = netid_match.group(1).strip()
            
            if 'author' in content_lower or 'name' in content_lower:
                name_match = re.search(r'(?:author|name|student\s+name)\s*[:=]\s*([^\n]+)', student_content, re.IGNORECASE)
                if name_match:
                    student_name = name_match.group(1).strip()
            
            results.append({
                'filename': filename,
                'student_name': student_name,
                'student_id': student_id,
                'score': result.get('score', 0),
                'feedback': result.get('feedback', ''),
                'question_feedback': result.get('question_feedback', []),
                'deductions': result.get('deductions', []),
                'debug': debug_info
            })
        
        return jsonify({'results': results})
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500


if __name__ == '__main__':
    # For production, set debug=False
    # For development, set debug=True
    debug_mode = os.getenv('FLASK_DEBUG', 'False').lower() == 'true'
    app.run(debug=debug_mode, host='0.0.0.0', port=4040)
